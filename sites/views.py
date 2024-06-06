import os
import uuid
import math
from datetime import datetime
import base64
import time
import requests
from pptx import Presentation
from django.conf import settings
from django.db.models import Q
from django.http import HttpResponse, JsonResponse
from rest_framework.decorators import api_view, permission_classes, parser_classes
from rest_framework.permissions import IsAuthenticated, IsAdminUser
from rest_framework.response import Response
from rest_framework import status
from rest_framework.parsers import MultiPartParser
from pptx.enum.shapes import MSO_SHAPE_TYPE
from django.core.exceptions import ObjectDoesNotExist
import requests
import datetime
import time
import pytz
import json


from .models import Site, Booking, SiteImages, SitePricing, GoogleStats, IndiaGrid, StateEntities, GoogleTrafficStats
from .serializers import SiteSerializer, BookingSerializer, ImagePathsSerializer, SitePricingSerializer, GooglePlacesSerializer, IndiaGridSerializer, StateEntitiesSerializer, GoogleTrafficSerializer

import googlemaps



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_india_grid(request):
    grid = IndiaGrid.objects.all()
    serializer = IndiaGridSerializer(grid, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def post_india_grid(request):
    data = request.data
    serializer = IndiaGridSerializer(data = data)
    if serializer.is_valid():
        serializer.save()
        # process_data(data)
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_state_entities(request):
    entity = StateEntities.objects.all()
    serializer = StateEntitiesSerializer(entity, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)


@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def post_state_entities(request):
    data = request.data
    serializer = StateEntitiesSerializer(data = data)
    if serializer.is_valid():
        serializer.save()
        # process_data(data)
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_site_detail(request, pk):
    site = Site.objects.get(pk=pk)
    serializer = SiteSerializer(site)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)




@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_sites_list(request):
    site = Site.objects.all()
    serializer = SiteSerializer(site, many=True)
    print(serializer)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def post_new_site(request):
    data = request.data
    data["siteTag"] = str(uuid.uuid4())
    if float(data["site_lat"]) < 8.4 or float(data["site_lat"]) > 37.6 or float(data["site_long"]) > 97.3 or float(data["site_long"]) < 68.7:
        Response({"success":False, "data":"Latitude and/or Longitude out of bounds"}, status=status.HTTP_400_BAD_REQUEST)
    data["grid"] = str(math.ceil((float(data["site_lat"])-8.4)/.01))+"_"+str(math.ceil((float(data["site_long"])-68.7)/.01))
    serializer = SiteSerializer(data = data)
    if serializer.is_valid():
        serializer.save()
        # process_data(data)
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def site_search_view(request):
    searchText = request.query_params.get('searchText')
    print("searchText", searchText)
    queryset = Site.objects.filter(
        Q(location__icontains=searchText) | Q(priorityArea__icontains=searchText) | Q(district__icontains=searchText)
    )
    serializer = SiteSerializer(queryset, many=True)

    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_site_price(request):
    siteTag = request.query_params.get('siteTag')
    try:
        pricing = SitePricing.objects.get(site=siteTag)
    except ObjectDoesNotExist:
        return Response({"success":False, "data":"No Pricing Found"}, status=status.HTTP_400_BAD_REQUEST)    
    serializer = SitePricingSerializer(pricing)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])    
def post_site_pricing(request):
    data = request.data
    
    serializer = SitePricingSerializer(data=data)
    
    if serializer.is_valid():
        serializer.save()
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    print(serializer.errors, "errors")
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)




@api_view(['PUT'])
@permission_classes([IsAuthenticated, IsAdminUser])
def edit_site_pricing(request):
    data = request.data
    try:
        pricing = SitePricing.objects.get(site=data["site"])
    except ObjectDoesNotExist:
        return Response({"success":False, "data":"No Pricing Found"}, status=status.HTTP_400_BAD_REQUEST)
    serializer = SitePricingSerializer(pricing, data)

    if serializer.is_valid():
        serializer.save()
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_202_ACCEPTED)
    print(serializer.errors, "errors")
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
@parser_classes([MultiPartParser])
def images_from_ppt(request):
    print(request, "Presentation Object")
    ppt_file = request.FILES.get('ppt_file')
    if ppt_file is None:
        return Response({'error': 'No file uploaded'}, status=400)
    
    images_folder = os.path.join(settings.STATIC_ROOT, 'ppt_images')
    
    if not os.path.exists(images_folder):
        os.makedirs(images_folder)

    presentation = Presentation(ppt_file)
    
    sitesUploaded = []
    slideNotUploaded = []
    index = 0
    for slide in presentation.slides:
        image = None
        text = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
            
            if shape.has_text_frame:
                if not text:
                    text = shape.text_frame.text.strip()
                
        if image and text:
            try:
                image_filename = os.path.join(images_folder, text+".jpeg")
                image_bytes = image.blob
                with open(image_filename, 'wb') as f:
                    f.write(image_bytes)
                sitesUploaded.append({"siteName": text, "path": image_filename})
            except:
                slideNotUploaded.append("Slide " + str(index) + " not uploaded due to name error")

        elif not image or not text:
            if not image:
                slideNotUploaded.append("Slide " + str(index) + " not uploaded as image not there")
            if not text: 
                slideNotUploaded.append("Slide " + str(index) + " not uploaded as text not there")
        
        index = index+1

    serializer = ImagePathsSerializer(data=sitesUploaded, many=True)
    if serializer.is_valid():
        serializer.save()
        uploaded_ids = [image['id'] for image in serializer.data]
        return Response({"success":True, "data":{'images': uploaded_ids, 'notUploaded':slideNotUploaded}}, status=status.HTTP_201_CREATED)
    else:
        return Response({"success":True, "data":'some error happened'}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def map_image_site(request):

    imageId = request.query_params.get('imageId')
    siteTag = request.query_params.get('siteTag')

    imageInstance = SiteImages.objects.get(pk=imageId)
    imageInstance.site = Site.objects.get(pk=siteTag)
    imageInstance.save()
    return Response({"success":True, "data":f'{siteTag} updated successfully'}, status=status.HTTP_200_OK)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
@parser_classes([MultiPartParser])
def get_site_images(request):
    siteTag = request.query_params.get('siteTag')
    print(siteTag)

    queryset = SiteImages.objects.filter(site=siteTag)[0]
    print(queryset)
    serializer = ImagePathsSerializer(queryset)
    image_data = None
    print(serializer.data["path"])
    with open(serializer.data["path"], 'rb') as f:
        image_data = f.read()
    response = HttpResponse(content_type='image/jpeg')
    response.write(image_data)
    return response



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
@parser_classes([MultiPartParser])
def upload_image_data(request):
    image_file = request.FILES.get('image')
    siteName = request.data.get('siteName')
    siteTag = request.data.get('siteTag')
    
    if image_file is None:
        return Response({'error': 'No file uploaded'}, status=400)

    images_folder = os.path.join(settings.STATIC_ROOT, 'ppt_images')
    
    if not os.path.exists(images_folder):
        os.makedirs(images_folder)

    now = datetime.now().strftime("%Y%m%dH%M%S")
    image_filename = os.path.join(images_folder, siteTag + now +".jpeg")
    image_bytes = image_file.read()
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)
        print("File is written")
                
    serializer = ImagePathsSerializer(data={"siteName": siteName, "path": image_filename, "site":siteTag})
    if serializer.is_valid():
        serializer.save()
        return Response({"success":True, "data":{"siteName": siteName, "path": image_filename, "site":siteTag}}, status=status.HTTP_201_CREATED)
    else:
        return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
@parser_classes([MultiPartParser])
def get_image_id(request, pk):
    queryset = SiteImages.objects.get(pk = pk)
    serializer = ImagePathsSerializer(queryset)
    image_data = None
    print(serializer.data)
    with open(serializer.data["path"], 'rb') as f:
        image_data = f.read()
    
    print(image_data)
    return JsonResponse({'image':base64.b64encode(image_data).decode('utf-8'), 'otherinfo':serializer.data})



@api_view(['GET'])
@permission_classes([IsAuthenticated, IsAdminUser])
def get_all_bookings(request):
    booking = Booking.objects.all()
    serializer = BookingSerializer(booking, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_self_bookings(request):
    company = request.user.company.code
    booking = Booking.objects.filter(assigned_to=company)
    serializer = BookingSerializer(booking, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def post_booking(request):
    data = request.data
    serializer = BookingSerializer(data=data)
    if serializer.is_valid():
        serializer.save()
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    return Response({"success":True, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)



@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def get_bookings_by_customer(request):
    searchText = request.query_params.get('customerID')
    queryset = Booking.objects.filter(assigned_to=searchText)
    serializer = BookingSerializer(queryset, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_nearby_assets_db(request):
    grid = request.query_params.get('grid')
    radius = request.query_params.get('radius', 1000)
    dataPoints = GoogleStats.objects.filter(Q(grid=grid))
    serializer = GooglePlacesSerializer(dataPoints, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)



@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_nearby_assets(request):
    grid = request.query_params.get('grid')
    radius = request.query_params.get('radius', 1000)
    keyword = request.query_params.get('keyword')

    
    dataPoints = GoogleStats.objects.filter(Q(grid=grid) & Q(type=keyword))
    if len(dataPoints) > 0:
        serializer = GooglePlacesSerializer(dataPoints, many=True)
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)

    [latSpan, longSpan] = grid.split("_")
    lat_min =  8.4
    long_min = 68.7
    location = (lat_min+float(latSpan)*0.01-0.005, long_min+float(longSpan)*0.01-0.005)
    
    gmaps = googlemaps.Client("AIzaSyAmnjoE7eFekledKAqF8ctrMoNBk3w0Xm8")
    print(gmaps, "GMAPS REACHED HERE")
    places = gmaps.places_nearby(
        location=location,
        radius=radius,
        keyword=keyword
    )
    
    results = []

    for place in places['results']:
        result = {
            'grid':grid,
            'type':keyword,
            'name': place['name'],
            'address': place['vicinity'],
            'status':place['business_status'],
            'latitude': place['geometry']['location']['lat'],
            'longitude': place['geometry']['location']['lng'],
            'place_id': place['place_id']
        }
        serializer = GooglePlacesSerializer(data = result)
        if serializer.is_valid():        
            serializer.save()
            results.append(result)


    while True:
        time.sleep(2)
        if "next_page_token" not in places:
            break
        else:
            next_page_token = places["next_page_token"]  
        print(next_page_token, "next page token")
        places = gmaps.places_nearby(
            location=location,
            radius=radius,
            keyword=keyword,
            page_token=next_page_token
        )

        print(places, "next Page Data")
        for place in places['results']:
            result = {
                'grid':grid,
                'type':keyword,
                'name': place['name'],
                'address': place['vicinity'],
                'status':place['business_status'],
                'latitude': place['geometry']['location']['lat'],
                'longitude': place['geometry']['location']['lng'],
                'place_id': place['place_id']
            }
            serializer = GooglePlacesSerializer(data = result)
            if serializer.is_valid():        
                serializer.save()
                results.append(result)

    return Response({"success":True, "data":results}, status=status.HTTP_200_OK)
    
@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_site_traffic(request):
    siteTag = request.query_params.get('siteTag')
    dataPoints = GoogleTrafficStats.objects.filter(Q(site=siteTag))
    serializer = GoogleTrafficSerializer(dataPoints, many=True)
    return Response({"success":True, "data":serializer.data}, status=status.HTTP_200_OK)


@api_view(['POST'])
@permission_classes([IsAuthenticated, IsAdminUser])
def post_new_traffic_request(request):
    data = request.data
    serializer = GoogleTrafficSerializer(data = data)
    if serializer.is_valid():
        serializer.save()
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_first_distance(request):
    origin = request.query_params.get('startCoord')
    destination = request.query_params.get('siteCoord')
    url = "https://maps.googleapis.com/maps/api/distancematrix/json"
    params = {
        "departure_time":"now",
        "origins": origin,
        "destinations": destination,
        "key": "AIzaSyAmnjoE7eFekledKAqF8ctrMoNBk3w0Xm8"
    }

    response = requests.get(url, params=params)
    data = response.json()
    return Response({"success":True, "data":data}, status=status.HTTP_200_OK)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def start_traffic_collection(request):
    site = request.query_params.get('siteTag')
    origin = request.query_params.get('startCoord')
    destination = request.query_params.get('siteCoord')
    url = "https://maps.googleapis.com/maps/api/distancematrix/json"
    params = {
        "departure_time":"now",
        "origins": origin,
        "destinations": destination,
        "key": "AIzaSyAmnjoE7eFekledKAqF8ctrMoNBk3w0Xm8"
    }

    response = requests.get(url, params=params)
    data_parsed = response.json()
    print(data_parsed, "data_parsed")
    data = {}
    data["site"] = site
    data["latitude_1"] = origin.split(',')[0]
    data["longitude_1"] = origin.split(',')[1]
    data["latitude_2"] = destination.split(',')[0]
    data["longitude_2"] = destination.split(',')[1]
    data["distance"] = data_parsed["rows"][0]["elements"][0]["distance"]["value"]
    data["average_time"] = data_parsed["rows"][0]["elements"][0]["duration"]["value"]
    data["traffic_time"] = data_parsed["rows"][0]["elements"][0]["duration_in_traffic"]["value"]

    serializer = GoogleTrafficSerializer(data = data)
    if serializer.is_valid():
        serializer.save()
        return Response({"success":True, "data":serializer.data}, status=status.HTTP_201_CREATED)
    return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)








@api_view(['POST'])
@permission_classes([IsAuthenticated])
def get_traffic_count_collection(request):
    days_of_week = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
    markers = json.loads(request.body.decode('utf-8'))
    #markers = request.body
    print(markers, "get_traffic_count_collection")
    outerCordinates = markers["outboundPoint"]
    siteCoordinates = markers["siteCoord"]
    site = markers["site"]
    print(markers, "get_traffic_count_collection")
    url = "https://maps.googleapis.com/maps/api/distancematrix/json"

    tomorrow = datetime.date.today() + datetime.timedelta(days=1)
    days = [tomorrow + datetime.timedelta(days=i) for i in range(7)]
    hours = {
             "morningPeak":[7,8,9,10,11],
             "noonHours": [12, 14, 16], 
             "eveningPeak":[17,18,19,20,21],
             "nightHours":[22, 23]
            }
    LOCAL_TIMEZONE = pytz.timezone('Asia/Kolkata')
    for outerMarker in outerCordinates:
        for day in days:
            for key, value in hours.items():
                for hour in value:
                    local_time = LOCAL_TIMEZONE.localize(datetime.datetime.combine(day, datetime.time(hour, 0)))
                    dayofWeek = days_of_week[local_time.weekday()]
                    utc_time = local_time.astimezone(pytz.utc)
                    departure_time = int(utc_time.timestamp())
                    url = 'https://maps.googleapis.com/maps/api/distancematrix/json'
                    params = {
                        'origins': f"{siteCoordinates["latitude"]},{siteCoordinates["longitude"]}",
                        'destinations': f"{outerMarker["lat"]},{outerMarker["lng"]}",
                        'departure_time': departure_time,
                        'key': "AIzaSyAmnjoE7eFekledKAqF8ctrMoNBk3w0Xm8",
                        'traffic_model': 'best_guess'
                    }
                    
                    response = requests.get(url, params=params)
                    data_parsed = response.json()
                    print(params, "params")
                    print(data_parsed, "params")
                    if response.status_code == 200 and data_parsed["rows"][0]["elements"][0]["status"] != "ZERO_RESULTS":
                        
                        print(data_parsed, "data_parsed")
                        data = {}
                        data["site"] = site
                        data["latitude_1"] = siteCoordinates["latitude"]
                        data["longitude_1"] = siteCoordinates["latitude"]
                        data["latitude_2"] = outerMarker["lat"]
                        data["longitude_2"] = outerMarker["lng"]
                        data["day"] = dayofWeek
                        data["daySection"] = key                        
                        data["distance"] = data_parsed["rows"][0]["elements"][0]["distance"]["value"]                
                        data["traffic_time"] = data_parsed["rows"][0]["elements"][0]["duration_in_traffic"]["value"]

                        serializer = GoogleTrafficSerializer(data = data)
                        if serializer.is_valid():
                            serializer.save()
                    time.sleep(1)

    return Response({"success":True, "data":"Created Successfully"}, status=status.HTTP_201_CREATED)
    # return Response({"success":False, "data":serializer.errors}, status=status.HTTP_400_BAD_REQUEST)

